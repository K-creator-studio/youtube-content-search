
import os
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
import nltk
from datetime import timedelta, datetime
import mysql.connector
from document_reader import get_default_file_path
import re
import argparse

nltk.download('punkt')

class SemanticTranscriptMatcher:
    def __init__(self, similarity_threshold=0.6, chunk_size=5, min_chunk_length=20):
        self.model = SentenceTransformer('paraphrase-MiniLM-L3-v2', device='cpu')
        self.similarity_threshold = similarity_threshold
        self.chunk_size = chunk_size
        self.min_chunk_length = min_chunk_length

    def extract_text(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            return ' '.join(
                run.text
                for slide in prs.slides
                for shape in slide.shapes
                if shape.has_text_frame
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
            )
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            raise ValueError(f"Unsupported file type: {file_path}")

    def parse_transcript_happysearch_style(self, file_path):
        entries = []
        with open(file_path, 'r', encoding='utf-8') as f:
            current_entry = {'timestamp': None, 'text': ''}
            for line in f:
                line = line.strip()
                if not line:
                    continue
                if '-->' in line:
                    if current_entry['text']:
                        entries.append(current_entry)
                    current_entry = {'timestamp': line, 'text': ''}
                else:
                    current_entry['text'] += (' ' + line) if current_entry['text'] else line
            if current_entry['text']:
                entries.append(current_entry)
        return entries

    def split_into_chunks(self, text):
        sentences = nltk.sent_tokenize(text)
        chunks = []
        for i in range(0, len(sentences), self.chunk_size):
            chunk = ' '.join(sentences[i:i + self.chunk_size])
            if len(chunk) >= self.min_chunk_length:
                chunks.append(chunk)
        return chunks

    def build_faiss_index(self, embeddings):
        dim = embeddings.shape[1]
        index = faiss.IndexFlatIP(dim)
        index.add(embeddings)
        return index

    @staticmethod
    def timestamp_to_seconds(ts):
        parts = ts.split(':')
        seconds = float(parts[2]) if '.' in parts[2] else int(parts[2])
        total_seconds = int(parts[0]) * 3600 + int(parts[1]) * 60 + seconds
        return total_seconds

    @staticmethod
    def parse_interval(interval_str):
        start_str, end_str = [s.strip() for s in interval_str.split('-->')]
        return (SemanticTranscriptMatcher.timestamp_to_seconds(start_str),
                SemanticTranscriptMatcher.timestamp_to_seconds(end_str))

    @staticmethod
    def seconds_to_timestamp(seconds):
        td = timedelta(seconds=seconds)
        total_seconds = int(td.total_seconds())
        hours = total_seconds // 3600
        minutes = (total_seconds % 3600) // 60
        secs = total_seconds % 60
        return f"{hours:02d}:{minutes:02d}:{secs:02d}"

    @staticmethod
    def merge_intervals(intervals):
        if not intervals:
            return []
        intervals.sort(key=lambda x: x[0])
        merged = [intervals[0]]
        for current in intervals[1:]:
            last = merged[-1]
            if current[0] <= last[1] + 1:
                merged[-1] = (last[0], max(last[1], current[1]))
            else:
                merged.append(current)
        return merged

    def match_document_with_transcripts(self, document_path, transcripts_folder, top_k=3):
        doc_text = self.extract_text(document_path)
        doc_chunks = self.split_into_chunks(doc_text)
        doc_embeddings = self.model.encode(doc_chunks, batch_size=32, show_progress_bar=False, normalize_embeddings=True).astype('float32')

        index = self.build_faiss_index(doc_embeddings)

        results = []
        for file_name in os.listdir(transcripts_folder):
            if not file_name.endswith('.txt'):
                continue
            transcript_path = os.path.join(transcripts_folder, file_name)
            transcript_entries = self.parse_transcript_happysearch_style(transcript_path)

            for entry in transcript_entries:
                text = entry['text']
                if not text.strip():
                    continue
                entry_embedding = self.model.encode([text], normalize_embeddings=True).astype('float32')
                distances, indices = index.search(entry_embedding, top_k)

                for dist, idx in zip(distances[0], indices[0]):
                    if dist >= self.similarity_threshold:
                        results.append({
                            'transcript_file': file_name,
                            'timestamp': entry['timestamp'],
                            'transcript_chunk': text,
                            'document_chunk': doc_chunks[idx],
                            'similarity': dist
                        })

        results.sort(key=lambda x: x['similarity'], reverse=True)
        return results

    def print_grouped_matches(self, matches):
        matches_by_file = {}
        for m in matches:
            matches_by_file.setdefault(m['transcript_file'], []).append(m)

        for file, file_matches in matches_by_file.items():
            intervals = [self.parse_interval(m['timestamp']) for m in file_matches]
            merged_intervals = self.merge_intervals(intervals)
            merged_intervals_str = [f"{self.seconds_to_timestamp(s)} --> {self.seconds_to_timestamp(e)}" for s, e in merged_intervals]

            overall_start = min(s for s, e in merged_intervals)
            overall_end = max(e for s, e in merged_intervals)
            overall_range_str = f"{self.seconds_to_timestamp(overall_start)} --> {self.seconds_to_timestamp(overall_end)}"

            top_match = max(file_matches, key=lambda x: x['similarity'])

            print(f"\nTranscript File: {file}")
            print(f"Timestamp: {overall_range_str}")
            for interval_str in merged_intervals_str:
                print(interval_str)
            print(f"Similarity: {top_match['similarity']:.2f}")
            print(f"Document Chunk: {top_match['document_chunk'][:200]}...")
            print(f"Transcript Chunk: {top_match['transcript_chunk'][:200]}...")

    def save_matches_to_file(self, matches, filename="matching_transcripts.txt"):
        if not matches:
            print("No matches to save.")
            return
        with open(filename, "w", encoding="utf-8") as f:
            matches_by_file = {}
            for m in matches:
                matches_by_file.setdefault(m['transcript_file'], []).append(m)

            for file, file_matches in matches_by_file.items():
                f.write(f"Transcript File: {file}\n")
                intervals = [self.parse_interval(m['timestamp']) for m in file_matches]
                merged_intervals = self.merge_intervals(intervals)
                overall_start = min(s for s, e in merged_intervals)
                overall_end = max(e for s, e in merged_intervals)
                overall_range_str = f"{self.seconds_to_timestamp(overall_start)} --> {self.seconds_to_timestamp(overall_end)}"
                f.write(f"Timestamp: {overall_range_str}\n")
                for interval_str in [f"{self.seconds_to_timestamp(s)} --> {self.seconds_to_timestamp(e)}" for s, e in merged_intervals]:
                    f.write(interval_str + "\n")
                top_match = max(file_matches, key=lambda x: x['similarity'])
                f.write(f"Similarity: {top_match['similarity']:.2f}\n")
                f.write(f"Document Chunk: {top_match['document_chunk'][:500]}...\n")
                f.write(f"Transcript Chunk: {top_match['transcript_chunk'][:500]}...\n")
                f.write("\n" + "-"*80 + "\n\n")
        print(f"Matching results saved to {filename}")

    def save_top_match_to_db(self, matches, db_config, course_id=None, lecture_id=None, video_title=None, video_url=None):
        if not matches:
            print("No matches to save to database.")
            return

        matches_by_file = {}
        for m in matches:
            matches_by_file.setdefault(m['transcript_file'], []).append(m)

        top_file = max(matches_by_file.keys(), 
                      key=lambda f: max(m['similarity'] for m in matches_by_file[f]))
        file_matches = matches_by_file[top_file]

        intervals = [self.parse_interval(m['timestamp']) for m in file_matches]
        merged_intervals = self.merge_intervals(intervals)
        
        start_time = self.seconds_to_timestamp(min(s for s, e in merged_intervals))
        end_time = self.seconds_to_timestamp(max(e for s, e in merged_intervals))
        
        top_similarity = max(m['similarity'] for m in file_matches)

        # Default values if not provided
        if not video_title:
            video_title = f"Matched Segment from {top_file}"
        if not video_url:
            video_url = "https://www.youtube.com/watch?v=matched_segment"

        try:
            conn = mysql.connector.connect(**db_config)
            cursor = conn.cursor()

            # If course_id and lecture_id not provided, try to extract from document
            if course_id is None or lecture_id is None:
                doc_name = os.path.basename(get_default_file_path())
                match = re.search(r'L(\d+)', doc_name)
                if match:
                    lecture_id = int(match.group(1))
                    course_id = 100 + lecture_id  # Example logic - adjust as needed

            insert_query = """
            INSERT INTO video_links 
            (course_id, lecture_id, video_title, video_url, start_time, end_time, added_at)
            VALUES (%s, %s, %s, %s, %s, %s, %s)
            """
            
            current_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            cursor.execute(insert_query, (
                course_id,
                lecture_id,
                video_title,
                video_url,
                start_time,
                end_time,
                current_time
            ))
            
            conn.commit()
            print(f"Successfully saved top match to database: {video_title} ({start_time} to {end_time})")
            
        except mysql.connector.Error as err:
            print(f"Error saving to database: {err}")
        finally:
            if conn.is_connected():
                cursor.close()
                conn.close()

    def get_similarity_whole_transcript(self, document_path, transcript_path):
        doc_text = self.extract_text(document_path)
        transcript_text = self.extract_text(transcript_path)

        doc_embedding = self.model.encode([doc_text], normalize_embeddings=True).astype('float32')
        transcript_embedding = self.model.encode([transcript_text], normalize_embeddings=True).astype('float32')

        similarity = float(np.dot(doc_embedding, transcript_embedding.T))
        return similarity


def get_db_config():
    """Helper function to get database configuration"""
    return {
        'host': 'localhost',
        'user': 'root',
        'password': '',
        'database': 'youtube_search_app'
    }

def parse_arguments():
    parser = argparse.ArgumentParser(description='Match document content with video transcripts')
    parser.add_argument('--course_id', type=int, help='Course ID for database')
    parser.add_argument('--lecture_id', type=int, help='Lecture ID for database')
    parser.add_argument('--video_title', type=str, help='Title for the video segment')
    parser.add_argument('--video_url', type=str, help='URL of the video')
    return parser.parse_args()

if __name__ == "__main__":
    args = parse_arguments()
    matcher = SemanticTranscriptMatcher()

    document_path = get_default_file_path()
    transcripts_folder = 'transcripts'

    matches = matcher.match_document_with_transcripts(document_path, transcripts_folder)

    print(f"Top semantic matches grouped by transcript file:")
    matcher.print_grouped_matches(matches)

    matcher.save_matches_to_file(matches, "matching_transcripts.txt")
    
    # Save top match to database
    matcher.save_top_match_to_db(
        matches,
        db_config=get_db_config(),
        course_id=args.course_id,
        lecture_id=args.lecture_id,
        video_title=args.video_title,
        video_url=args.video_url
    )

    print("\nOverall similarity between document and each full transcript file:")
    for file_name in os.listdir(transcripts_folder):
        if not file_name.endswith('.txt'):
            continue
        transcript_path = os.path.join(transcripts_folder, file_name)
        overall_sim = matcher.get_similarity_whole_transcript(document_path, transcript_path)
        print(f"{file_name}: {overall_sim:.4f}")
