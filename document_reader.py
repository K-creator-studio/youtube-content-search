

import os
os.environ["HF_HOME"] = "D:/KIRAN/fyp/PROJECT/.cache/huggingface"
import re
import string
import fitz  # PyMuPDF
import easyocr
import nltk
from pptx import Presentation
from docx import Document
from transformers import pipeline
from rake_nltk import Rake
from keybert import KeyBERT
from nltk.corpus import stopwords
import spacy
from collections import Counter
from sklearn.feature_extraction.text import TfidfVectorizer


# Set NLTK data directory to D drive
nltk_data_dir = os.path.join("D:/", "nltk_data")
os.makedirs(nltk_data_dir, exist_ok=True)  # Create directory if it doesn't exist
nltk.data.path.append(nltk_data_dir)  # Add to NLTK's data path

# -------------------------------
# Setup and Initialization
# -------------------------------
# Download required NLTK resources
# Replace your current NLTK download calls with:
nltk.download('punkt', download_dir=nltk_data_dir, quiet=True)
nltk.download('stopwords', download_dir=nltk_data_dir, quiet=True)
nltk.download('wordnet', download_dir=nltk_data_dir, quiet=True)
nltk.download('averaged_perceptron_tagger', download_dir=nltk_data_dir, quiet=True)

# Initialize models
reader = easyocr.Reader(['en'], model_storage_directory='D:/KIRAN/fyp/EasyOCR_cache')
  # For OCR on images
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
kw_model = KeyBERT()

# Load spaCy model
try:
    nlp = spacy.load("en_core_web_md")  # Medium-sized model with word vectors
except:
    # If model not installed, download it
    import subprocess
    subprocess.run(["python", "-m", "spacy", "download", "en_core_web_md"])
    nlp = spacy.load("en_core_web_md")

# Custom stopwords list combining NLTK and domain-specific terms
CUSTOM_STOPWORDS = set(stopwords.words('english') +
                      ['example', 'etc', 'ie', 'eg', 'figure', 'table', 'section',
                       'chapter', 'page', 'content', 'slide', 'presentation', 'adapted',
                       'using', 'used', 'use', 'uses', 'one', 'two', 'three', 'first',
                       'second', 'third', 'also', 'may', 'many', 'much', 'however',
                       'thus', 'therefore', 'hence', 'although', 'though', 'despite',
                       'nevertheless', 'nonetheless', 'yet', 'still', 'otherwise',
                       'accordingly', 'consequently', 'example', 'examples', 'shown',
                       'shows', 'show', 'showing', 'shown', 'see', 'sees', 'seen',
                       'seeing', 'look', 'looks', 'looking', 'looked', 'appear',
                       'appears', 'appeared', 'appearing', 'seem', 'seems', 'seemed',
                       'seeming', 'like', 'likes', 'liked', 'liking'])

# -------------------------------
# Text Extraction Functions
# -------------------------------
def extract_text_from_pdf(pdf_path):
    """Extract text from PDF files"""
    text = ""
    try:
        doc = fitz.open(pdf_path)
        for page in doc:
            text += page.get_text()
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PDF: {str(e)}")
        return ""

def extract_text_from_ppt(ppt_path):
    """Extract text from PowerPoint presentations"""
    text = ""
    try:
        prs = Presentation(ppt_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {str(e)}")
        return ""

def extract_text_from_word(docx_path):
    """Extract text from Word documents"""
    text = ""
    try:
        doc = Document(docx_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from Word document: {str(e)}")
        return ""

def extract_text_from_image(image_path):
    """Extract text from images using OCR"""
    try:
        result = reader.readtext(image_path, detail=0)
        return " ".join(result)
    except Exception as e:
        print(f"Error extracting text from image: {str(e)}")
        return ""

def extract_text(file_path):
    """Extract text from various file types"""
    if not os.path.exists(file_path):
        return f"File not found: {file_path}"

    ext = os.path.splitext(file_path)[-1].lower()

    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in [".ppt", ".pptx"]:
        return extract_text_from_ppt(file_path)
    elif ext in [".doc", ".docx"]:
        return extract_text_from_word(file_path)
    elif ext in [".png", ".jpg", ".jpeg"]:
        return extract_text_from_image(file_path)
    elif ext in [".txt", ".md", ".csv"]:
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            return f"Error reading file: {str(e)}"
    else:
        return f"Unsupported file format: {ext}"

# -------------------------------
# Text Preprocessing
# -------------------------------
def preprocess_text(text):
    """Clean and preprocess text for keyword extraction"""
    if not text:
        return ""

    # Convert to lowercase and remove extra whitespace
    text = re.sub(r'\s+', ' ', text.lower()).strip()

    # Remove URLs
    text = re.sub(r'http\S+|www\S+|https\S+', '', text)

    # Remove email addresses
    text = re.sub(r'\S*@\S*\s?', '', text)

    # Remove special characters but keep spaces between words
    text = re.sub(r'[^\w\s]', ' ', text)

    # Remove extra spaces
    text = re.sub(r'\s+', ' ', text).strip()

    return text

# -------------------------------
# Summarization Functions
# -------------------------------
def summarize_text(text, max_length=1000, min_length=50):
    """Summarize text using the BART model"""
    if not text or len(text) < 100:
        return text

    try:
        summary = summarizer(text, max_length=max_length, min_length=min_length, do_sample=False)
        return summary[0]['summary_text']
    except Exception as e:
        print(f"Error summarizing text: {str(e)}")
        # Fallback to first few sentences
        sentences = nltk.sent_tokenize(text)
        return " ".join(sentences[:3])

def summarize_large_text(text, chunk_size=1024, max_length=1000, min_length=50):
    """Summarize large text by breaking it into chunks"""
    if not text:
        return ""

    try:
        # Split text into chunks
        chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]

        # Summarize each chunk
        chunk_summaries = []
        for chunk in chunks:
            if len(chunk.strip()) > 100:  # Only summarize substantial chunks
                summary = summarizer(chunk, max_length=max_length, min_length=min_length, do_sample=False)
                chunk_summaries.append(summary[0]['summary_text'])
            else:
                chunk_summaries.append(chunk)

        # Combine chunk summaries
        combined_summary = " ".join(chunk_summaries)

        # If the combined summary is still too long, summarize it again
        if len(combined_summary) > chunk_size:
            return summarize_text(combined_summary, max_length=max_length*2, min_length=min_length)

        return combined_summary
    except Exception as e:
        print(f"Error summarizing large text: {str(e)}")
        # Fallback to first few sentences
        sentences = nltk.sent_tokenize(text)
        return " ".join(sentences[:5])

# -------------------------------
# Keyword Validation
# -------------------------------
def is_valid_keyword(phrase, min_length=3, max_length=50):
    """Enhanced validation for potential keywords"""
    if not phrase:
        return False

    phrase_clean = phrase.strip()
    lc_phrase = phrase_clean.lower()

    # Basic length and format checks
    if (len(phrase_clean) < min_length or
        len(phrase_clean) > max_length or
        lc_phrase in CUSTOM_STOPWORDS or
        all(char in string.punctuation for char in lc_phrase) or
        lc_phrase.replace(" ", "").isnumeric()):
        return False

    # Check for URLs, emails, file paths
    if (re.search(r'http|www|\.(com|org|net|edu)|@|\\|/', lc_phrase) or
        re.match(r'^[a-z]:\\', lc_phrase)):
        return False

    # Check for dates and times
    if re.search(r'\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|\d{1,2}:\d{2}', lc_phrase):
        return False

    # Check for citations and references
    if re.search(r'^\[\d+\]|^fig\.|^table|^et al', lc_phrase):
        return False

    # Ensure it has at least one meaningful word (not just numbers and symbols)
    if not re.search(r'[a-z]{3,}', lc_phrase):
        return False

    return True

# -------------------------------
# Keyword Extraction Methods
# -------------------------------
def extract_rake_keywords(text, top_n=20):
    """Extract keywords using RAKE algorithm"""
    if not text or len(text) < 10:
        return []

    try:
        rake = Rake(stopwords=list(CUSTOM_STOPWORDS), min_length=1, max_length=3)
        rake.extract_keywords_from_text(text)
        keywords = [phrase for _, phrase in rake.get_ranked_phrases_with_scores()
                   if is_valid_keyword(phrase)]
        return keywords[:top_n]
    except Exception as e:
        print(f"RAKE extraction error: {str(e)}")
        return []

def extract_keybert_keywords(text, top_n=20):
    """Extract keywords using KeyBERT (semantic extraction)"""
    if not text or len(text) < 50:
        return []

    try:
        keybert_results = kw_model.extract_keywords(
            text,
            keyphrase_ngram_range=(1, 2),
            stop_words=list(CUSTOM_STOPWORDS),
            use_mmr=True,  # Use Maximal Marginal Relevance for diversity
            diversity=0.5,
            top_n=top_n
        )
        keywords = [kw[0] for kw in keybert_results if is_valid_keyword(kw[0])]
        return keywords
    except Exception as e:
        print(f"KeyBERT extraction error: {str(e)}")
        return []

def extract_spacy_keywords(text, top_n=20):
    """Extract keywords using spaCy's NER and noun chunk extraction"""
    if not text or len(text) < 10:
        return []

    try:
        doc = nlp(text[:100000])  # Limit text size to avoid memory issues

        # Extract named entities
        entities = [ent.text.lower() for ent in doc.ents
                   if ent.label_ in ['ORG', 'PRODUCT', 'WORK_OF_ART', 'LAW', 'LANGUAGE', 'EVENT', 'FAC']]

        # Extract noun chunks (noun phrases)
        noun_chunks = [chunk.text.lower() for chunk in doc.noun_chunks
                      if len(chunk.text.split()) <= 3]  # Limit to 3 words max

        # Extract technical terms (adjective + noun patterns)
        tech_terms = []
        for token in doc:
            if token.pos_ == "NOUN":
                # Check for adjective + noun pattern
                if token.i > 0 and doc[token.i-1].pos_ == "ADJ":
                    term = doc[token.i-1].text + " " + token.text
                    tech_terms.append(term.lower())

        # Combine all candidates
        all_candidates = entities + noun_chunks + tech_terms

        # Filter and count occurrences
        filtered_candidates = [cand for cand in all_candidates if is_valid_keyword(cand)]
        counter = Counter(filtered_candidates)

        # Return top N most common
        return [item[0] for item in counter.most_common(top_n)]
    except Exception as e:
        print(f"spaCy extraction error: {str(e)}")
        return []

def extract_tfidf_keywords(text, top_n=20):
    """Extract keywords using TF-IDF with adaptive parameters based on document size"""
    if not text or len(text) < 10:
        return []

    try:
        # Create sentences for the TF-IDF context
        sentences = nltk.sent_tokenize(text)

        # Adjust parameters based on number of sentences
        if len(sentences) <= 5:
            # For very few sentences, use minimal constraints
            min_df = 1
            max_df = 1.0
        elif len(sentences) <= 10:
            min_df = 1
            max_df = 0.9
        else:
            min_df = 2
            max_df = 0.85

        # Create TF-IDF vectorizer with adaptive parameters
        vectorizer = TfidfVectorizer(
            ngram_range=(1, 2),  # Consider unigrams and bigrams
            stop_words=list(CUSTOM_STOPWORDS),
            max_df=max_df,
            min_df=min_df,
        )

        # If we don't have enough sentences, add chunks of the text
        if len(sentences) < 3:
            # Split text into chunks to create more "documents"
            words = text.split()
            chunk_size = max(30, len(words) // 5)  # Aim for at least 5 chunks
            chunks = [' '.join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]
            documents = sentences + chunks
        else:
            documents = sentences

        # Fit and transform the documents
        tfidf_matrix = vectorizer.fit_transform(documents)

        # Get feature names and their scores
        feature_names = vectorizer.get_feature_names_out()

        # Sum TF-IDF scores across all documents for each term
        tfidf_scores = tfidf_matrix.sum(axis=0).A1

        # Create a dictionary of terms and their scores
        term_scores = {feature_names[i]: tfidf_scores[i] for i in range(len(feature_names))}

        # Sort terms by score and filter valid keywords
        sorted_terms = sorted(term_scores.items(), key=lambda x: x[1], reverse=True)
        filtered_terms = [term for term, score in sorted_terms if is_valid_keyword(term)]

        return filtered_terms[:top_n]

    except ValueError as e:
        print(f"TF-IDF extraction error (trying fallback): {str(e)}")
        # Fallback if TF-IDF still fails
        try:
            # Use a simple frequency-based approach instead
            words = nltk.word_tokenize(text.lower())
            bigrams = list(nltk.bigrams(words))

            # Count word frequencies (unigrams)
            word_freq = Counter([w for w in words if len(w) > 2 and w not in CUSTOM_STOPWORDS])

            # Count bigram frequencies
            bigram_freq = Counter([' '.join(bg) for bg in bigrams])

            # Combine unigrams and bigrams, prioritizing bigrams
            combined_terms = []
            for term, count in bigram_freq.most_common(top_n * 2):
                if is_valid_keyword(term):
                    combined_terms.append(term)

            for term, count in word_freq.most_common(top_n * 2):
                if is_valid_keyword(term) and len(combined_terms) < top_n:
                    combined_terms.append(term)

            return combined_terms[:top_n]
        except Exception as e2:
            print(f"TF-IDF fallback error: {str(e2)}")
            return []
    except Exception as e:
        print(f"TF-IDF extraction error: {str(e)}")
        return []

def extract_frequency_keywords(text, top_n=20):
    """Simple frequency-based keyword extraction as fallback"""
    if not text:
        return []

    try:
        # Tokenize and clean
        words = nltk.word_tokenize(text.lower())
        words = [w for w in words if len(w) > 3 and w not in CUSTOM_STOPWORDS]

        # Get word frequencies
        word_freq = Counter(words)

        # Get bigram frequencies
        bigrams = list(nltk.bigrams(words))
        bigram_freq = Counter([' '.join(bg) for bg in bigrams])

        # Combine results
        keywords = []

        # Add top bigrams
        for term, _ in bigram_freq.most_common(top_n):
            if is_valid_keyword(term) and term not in keywords:
                keywords.append(term)

        # Add top unigrams
        for term, _ in word_freq.most_common(top_n * 2):
            if is_valid_keyword(term) and term not in keywords and len(keywords) < top_n:
                keywords.append(term)

        return keywords[:top_n]
    except Exception as e:
        print(f"Frequency extraction error: {str(e)}")
        return []

# -------------------------------
# Combined Keyword Extraction
# -------------------------------
def extract_refined_keywords(text, summary=None, top_n=15):
    """
    Extract refined keywords using multiple methods and combine results

    Args:
        text: The full document text
        summary: Optional summary text (if available)
        top_n: Number of keywords to return

    Returns:
        List of refined keywords
    """
    if not text:
        return []

    # Clean the text
    clean_text = preprocess_text(text)

    # If summary is provided, use it for keyword extraction as well
    if summary:
        clean_summary = preprocess_text(summary)
        # Give more weight to summary-based keywords
        text_to_process = clean_summary + " " + clean_summary + " " + clean_text
    else:
        text_to_process = clean_text

    # Initialize empty lists for all keyword methods
    rake_keywords = []
    keybert_keywords = []
    spacy_keywords = []
    tfidf_keywords = []
    frequency_keywords = []

    # Try each method independently and catch any errors
    # 1. RAKE keywords
    rake_keywords = extract_rake_keywords(text_to_process, top_n=top_n)

    # 2. KeyBERT keywords
    keybert_keywords = extract_keybert_keywords(text_to_process, top_n=top_n)

    # 3. spaCy keywords
    spacy_keywords = extract_spacy_keywords(text_to_process, top_n=top_n)

    # 4. TF-IDF keywords
    tfidf_keywords = extract_tfidf_keywords(text_to_process, top_n=top_n)

    # 5. Frequency keywords (as backup)
    frequency_keywords = extract_frequency_keywords(text_to_process, top_n=top_n)

    # Combine all keywords with priority weighting
    keyword_counter = Counter()

    # Give different weights to different methods
    for kw in keybert_keywords:
        keyword_counter[kw] += 3  # Higher weight for semantic relevance

    for kw in spacy_keywords:
        keyword_counter[kw] += 2  # Good weight for linguistic patterns

    for kw in rake_keywords:
        keyword_counter[kw] += 1  # Lower weight for statistical approach

    for kw in tfidf_keywords:
        keyword_counter[kw] += 2  # Good weight for document relevance

    for kw in frequency_keywords:
        keyword_counter[kw] += 1  # Lowest weight for simple frequency

    # If we have no keywords from any method, fall back to simple word frequency
    if not keyword_counter:
        return frequency_keywords

    # Get the most common keywords across all methods
    final_keywords = [kw for kw, _ in keyword_counter.most_common(top_n*2)]

    # Remove any duplicates that might have different capitalization or spacing
    seen = set()
    unique_keywords = []
    for kw in final_keywords:
        kw_norm = re.sub(r'\s+', ' ', kw.lower()).strip()
        if kw_norm not in seen and is_valid_keyword(kw_norm):
            seen.add(kw_norm)
            unique_keywords.append(kw_norm)

    return unique_keywords[:top_n]



def clean_title(title):
    """
    Clean the title by removing unwanted characters (timestamps, repetitions, etc.),
    excluding words less than three characters, and preserving the meaningful words.

    Args:
        title (str): The raw title.

    Returns:
        str: The cleaned and meaningful title text.
    """
    # Convert to lowercase for consistency
    title = title.lower()

    # Remove excessive punctuation (e.g., dashes, underscores, multiple spaces)
    title = re.sub(r'[-_]+', ' ', title)  # Replace dashes or underscores with a single space
    title = re.sub(r'[^\w\s]', '', title)  # Remove non-alphanumeric characters (except spaces)

    # Remove redundant words like 'pm', 'am', and other similar terms
    title = re.sub(r'\b(pm|am)\b', '', title)  # Remove 'pm' and 'am' words

    # Remove sequences of digits (timestamps, version numbers, etc.)
    title = re.sub(r'\d+', '', title)  # Remove all digits

    # Remove any words that are less than 3 characters
    title = ' '.join([word for word in title.split() if len(word) >= 3])

    # Remove extra spaces between words
    title = re.sub(r'\s+', ' ', title).strip()  # Remove extra spaces

    # Capitalize the first letter of each word for readability
    title = title.title()

    return title



# -------------------------------
# Main Processing Function
# -------------------------------
# def process_document(file_path, keyword_count=15):
#     """
#     Main function to process a document: extract text, summarize, and extract keywords

#     Args:
#         file_path: Path to the document file
#         keyword_count: Number of keywords to extract

#     Returns:
#         Tuple of (extracted_text, summary, keywords)
#     """
#     try:
#         # Extract text from the document
#         content = extract_text(file_path)

#         if not content or len(content) < 10:
#             return "No content extracted or content too short", None, None

#         # Generate summary if content is long enough
#         summary = None
#         if len(content) > 100:
#             try:
#                 summary = summarize_large_text(content) if len(content) > 1024 else summarize_text(content)
#             except Exception as e:
#                 print(f"Summarization error: {str(e)}")
#                 # Use the first few sentences as a fallback summary
#                 sentences = nltk.sent_tokenize(content)
#                 summary = " ".join(sentences[:3])
#         else:
#             summary = content

#         # Extract keywords using our enhanced method
#         try:
#             keywords = extract_refined_keywords(content, summary, top_n=keyword_count)
#         except Exception as e:
#             import traceback
#             print(f"Keyword extraction error: {str(e)}")
#             print(traceback.format_exc())
#             # Fallback to simple word frequency
#             words = nltk.word_tokenize(content.lower())
#             words = [w for w in words if len(w) > 3 and w not in CUSTOM_STOPWORDS and is_valid_keyword(w)]
#             keywords = [w for w, _ in Counter(words).most_common(keyword_count)]

#         return content[:5000], summary, keywords

#     except Exception as e:
#         import traceback
#         return f"Error: {str(e)}\n{traceback.format_exc()}", None, None

def get_default_file_path():
    return "L13-Artificial-Neural-Network-27122022-120614pm-28122022-035917pm-06122023-094345am-3--14052024-125256pm.pptx" 



def process_document(file_path, keyword_count=15):
    """
    Main function to process a document: extract text, summarize, and extract keywords.

    Args:
        file_path: Path to the document file.
        keyword_count: Number of keywords to extract.

    Returns:
        Tuple of (extracted_text, summary, content_keywords, title_keyword, combined_keywords).
    """
    try:
        # Extract the file name (title) from the file path
        file_name = os.path.basename(file_path)
        raw_title = os.path.splitext(file_name)[0]

        # Clean the title to extract meaningful text
        title_keyword = clean_title(raw_title)

        # Extract text from the document
        content = extract_text(file_path)

        if not content or len(content) < 10:
            return "No content extracted or content too short", None, None, None, None

        # Generate summary if content is long enough
        summary = None
        if len(content) > 100:
            try:
                summary = summarize_large_text(content) if len(content) > 1024 else summarize_text(content)
            except Exception as e:
                print(f"Summarization error: {str(e)}")
                # Use the first few sentences as a fallback summary
                sentences = nltk.sent_tokenize(content)
                summary = " ".join(sentences[:3])
        else:
            summary = content

        # Extract keywords from the document content
        try:
            content_keywords = extract_refined_keywords(content, summary, top_n=keyword_count)
        except Exception as e:
            import traceback
            print(f"Keyword extraction error: {str(e)}")
            print(traceback.format_exc())
            # Fallback to simple word frequency
            words = nltk.word_tokenize(content.lower())
            words = [w for w in words if len(w) > 3 and w not in CUSTOM_STOPWORDS and is_valid_keyword(w)]
            content_keywords = [w for w, _ in Counter(words).most_common(keyword_count)]

        # Combine content and title keywords
        combined_keywords = list(set([title_keyword] + content_keywords))  # Add title as one keyword

        # Ensure title_keyword is the first keyword in the combined list
        if title_keyword in combined_keywords:
            combined_keywords.remove(title_keyword)  # Remove duplicate if it exists later in the list
        combined_keywords.insert(0, title_keyword)  # Insert title_keyword at the start

        return content[:5000], summary, content_keywords, title_keyword, combined_keywords

    except Exception as e:
        import traceback
        return f"Error: {str(e)}\n{traceback.format_exc()}", None, None, None, None

# -------------------------------
# Example Usage
# -------------------------------
if __name__ == "__main__":
    file_path = get_default_file_path()  # Replace with your document path
    extracted_text, summary, content_keywords, title_keyword, combined_keywords = process_document(file_path)

    if extracted_text and not extracted_text.startswith("Error"):
        print("Extracted Text (first 200 chars):\n", extracted_text[:200] + "...")
    #     print("\nSummary:\n", summary)
    #     print("\nKeywords:\n", keywords)
    # else:
    #     print("Error:", extracted_text)

        print("\nContent Keywords:\n", content_keywords)
        print("\nTitle Keyword:\n", title_keyword)
        print("\nCombined Keywords:\n", combined_keywords)
    else:
        print("Error:", extracted_text)